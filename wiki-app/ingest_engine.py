"""
服务端 Ingest 引擎

职责：
  1. 文件解析     — PDF / PPTX / DOCX → 纯文本
  2. 领域配置读取 — schema/domain-config.xlsx → 结构化描述
  3. 现有 wiki 扫描 — 提取实体索引，供 LLM 做版本对比
  4. Prompt 构建  — 将以上信息组装为 LLM 可理解的指令
  5. LLM 流式调用 — 同步生成器，供 FastAPI SSE 端点消费
  6. 结果解析     — 从 LLM 输出中提取 JSON
  7. 写入 wiki    — commit 到实际文件系统
"""

from __future__ import annotations

import json
import re
import textwrap
from datetime import date
from pathlib import Path
from typing import Generator, Optional

import frontmatter as _fm

# ─── 文件解析 ────────────────────────────────────────────────────────────────

_MAX_TEXT_CHARS = 40_000   # 传入 LLM 的最大字符数


def extract_text(file_path: Path) -> str:
    """
    从上传文件中提取纯文本。
    支持：.pdf / .pptx / .docx / .txt / .md
    """
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(file_path)
    elif suffix == ".pptx":
        return _extract_pptx(file_path)
    elif suffix in (".docx",):
        return _extract_docx(file_path)
    elif suffix in (".txt", ".md"):
        return file_path.read_text(encoding="utf-8", errors="replace")
    else:
        raise ValueError(f"不支持的文件类型：{suffix}，请上传 PDF / PPTX / DOCX")


def _extract_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("请安装 PyMuPDF：pip install PyMuPDF")

    doc = fitz.open(str(path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(f"--- 第 {i+1} 页 ---\n{text}")
    return _truncate("\n\n".join(pages))


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("请安装 python-pptx：pip install python-pptx")

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        if parts:
            slides.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(parts))
    return _truncate("\n\n".join(slides))


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("请安装 python-docx：pip install python-docx")

    doc = Document(str(path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # 表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return _truncate("\n".join(parts))


def _truncate(text: str) -> str:
    if len(text) <= _MAX_TEXT_CHARS:
        return text
    half = _MAX_TEXT_CHARS // 2
    return text[:half] + "\n\n[……内容过长，已截断中间部分……]\n\n" + text[-half:]


# ─── 领域配置读取 ─────────────────────────────────────────────────────────────

def load_domain_config(schema_dir: Path) -> str:
    """
    读取 domain-config.xlsx，返回结构化的文本描述供 LLM 使用。
    若文件不存在则返回空字符串（不中断流程）。
    """
    xlsx_path = schema_dir / "domain-config.xlsx"
    if not xlsx_path.exists():
        return "（domain-config.xlsx 未找到，请按通用规则提取）"

    try:
        import openpyxl
    except ImportError:
        return "（openpyxl 未安装，无法读取领域配置）"

    wb = openpyxl.load_workbook(str(xlsx_path), read_only=True, data_only=True)
    sections: list[str] = []

    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        lines = [f"### Sheet: {sheet.title}"]
        for row in rows:
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))
        sections.append("\n".join(lines))

    wb.close()
    return "\n\n".join(sections)


# ─── 现有 wiki 实体扫描 ───────────────────────────────────────────────────────

_SKIP_FILES = {
    "log.md", "index.md", "schema-suggestions.md",
    "pending-clarifications.md", "README.md",
}


def scan_existing_entities(wiki_path: Path) -> list[dict]:
    """
    扫描 wiki/ 目录，提取所有实体的 frontmatter 摘要。
    返回轻量索引列表，供 LLM 做版本对比用。
    """
    if not wiki_path.exists():
        return []

    result = []
    for md_file in sorted(wiki_path.glob("**/*.md")):
        if md_file.name in _SKIP_FILES:
            continue
        try:
            post = _fm.loads(md_file.read_text(encoding="utf-8"))
            meta = post.metadata
            if not meta:
                continue
            result.append({
                "wiki_path": str(md_file.relative_to(wiki_path)),
                "entity_id": meta.get("entity_id", ""),
                "title": meta.get("title", md_file.stem),
                "type": meta.get("type", ""),
                "status": meta.get("status", "active"),
                "version": meta.get("version", "v1"),
            })
        except Exception:
            pass

    return result


# ─── Prompt 构建 ──────────────────────────────────────────────────────────────

_SYSTEM_PROMPT = textwrap.dedent("""
你是企业知识库（LLM Wiki Kit）的结构化摄入引擎。
你的唯一任务是将原始文档内容，转化为标准的 wiki 知识条目。

## 输出要求（严格遵守）
- 只输出一个合法的 JSON 对象，不添加任何说明文字
- JSON 必须放在 ```json ... ``` 代码块中
- 字符串中的换行用 \\n 表示
- frontmatter 字段值中不得包含裸冒号（使用引号包裹）

## Wiki 核心规范
- 条目格式：YAML Frontmatter + Markdown 正文，存于 wiki/ 子目录
- entity_id 格式：{类型缩写}-{年份}-{关键词}-{序号:02d}，全小写，用连字符
- status: active | outdated | archived（新建默认 active）
- version: v1、v2……（更新时递增）
- 若判断是对现有实体的更新：action="update"，旧字段值写入 history_block
- 若是全新实体：action="create"，history_block=null

## 冲突类型
- Type A：同字段不同来源数值不同 → 两值并存，加 [CONFLICT] 标记
- Type B：定性矛盾 → 暂不写入，输出到 conflicts 数组
- Type C：观点差异（非事实冲突） → 两者均写入，加 [PERSPECTIVE] 标记

## 自发现字段
对 Schema 未定义但有高价值的字段，放入 schema_suggestions 数组而非直接写入 frontmatter。
""").strip()


def build_ingest_prompt(
    filename: str,
    extracted_text: str,
    domain_config: str,
    existing_entities: list[dict],
    schema_core: str,
) -> tuple[str, str]:
    """返回 (system_prompt, user_prompt)"""

    # 现有实体索引（精简，避免超出 context）
    if existing_entities:
        entity_lines = ["| wiki_path | title | type | status | version |", "|---|---|---|---|---|"]
        for e in existing_entities[:60]:   # 最多展示 60 条
            entity_lines.append(
                f"| {e['wiki_path']} | {e['title']} | {e['type']} | {e['status']} | {e['version']} |"
            )
        entities_section = "\n".join(entity_lines)
    else:
        entities_section = "（暂无现有实体）"

    today = date.today().isoformat()

    user_prompt = textwrap.dedent(f"""
## 领域配置（domain-config.xlsx）

{domain_config}

---

## 现有知识库实体索引（供版本对比）

{entities_section}

---

## 待摄入文件

文件名：{filename}
摄入日期：{today}

内容：
{extracted_text}

---

## 输出 JSON 格式（严格按此结构）

```json
{{
  "action": "create",
  "wiki_path": "clients/xxx.md",
  "entity_id": "client-2024-xxx-01",
  "title": "项目或实体名称",
  "frontmatter": {{
    "title": "...",
    "type": "client",
    "entity_id": "client-2024-xxx-01",
    "status": "active",
    "version": "v1",
    "created_at": "{today}",
    "updated_at": "{today}",
    "sources": ["{filename}"]
  }},
  "body": "## [summary] 摘要\\n正文内容……",
  "history_block": null,
  "conflicts": [
    {{
      "id": "conflict-001",
      "type": "B",
      "field": "字段名",
      "existing_value": "现有值",
      "new_value": "新值",
      "recommendation": "建议说明"
    }}
  ],
  "schema_suggestions": [
    {{
      "field": "建议字段名",
      "example_value": "示例值",
      "reason": "发现价值说明",
      "confidence": "high"
    }}
  ],
  "related_updates": [
    {{
      "wiki_path": "industries/某行业.md",
      "append_text": "\\n- 新增关联：[[clients/xxx]]"
    }}
  ],
  "summary": "一句话摘要"
}}
```

请严格按此格式输出，不添加其他文字。
""").strip()

    return _SYSTEM_PROMPT, user_prompt


# ─── LLM 流式调用（同步生成器） ───────────────────────────────────────────────

def stream_llm_ingest(
    system_prompt: str,
    user_prompt: str,
    client,
    model: str,
) -> Generator[str, None, None]:
    """
    同步流式生成器，逐 token yield LLM 输出。
    由 FastAPI SSE 端点在线程池中消费。
    """
    stream = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.1,
        max_tokens=4000,
        stream=True,
    )
    for chunk in stream:
        delta = chunk.choices[0].delta.content or ""
        if delta:
            yield delta


# ─── 结果解析 ─────────────────────────────────────────────────────────────────

def parse_llm_output(raw: str) -> dict:
    """
    从 LLM 输出中提取并解析 JSON。
    支持带 ```json...``` 代码块或裸 JSON。
    """
    # 优先提取代码块
    m = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
    json_str = m.group(1) if m else raw.strip()

    # 尝试解析
    try:
        result = json.loads(json_str)
    except json.JSONDecodeError as e:
        # 尝试修复常见问题：尾部逗号
        cleaned = re.sub(r",\s*([}\]])", r"\1", json_str)
        try:
            result = json.loads(cleaned)
        except json.JSONDecodeError:
            raise ValueError(
                f"LLM 输出不是合法 JSON（{e}）。请重试，或检查文件内容是否过于复杂。"
            )

    # 基础字段验证
    required = ("action", "wiki_path", "frontmatter", "body")
    missing = [k for k in required if k not in result]
    if missing:
        raise ValueError(f"LLM 输出缺少必要字段：{missing}")

    # 设置默认值
    result.setdefault("conflicts", [])
    result.setdefault("schema_suggestions", [])
    result.setdefault("related_updates", [])
    result.setdefault("history_block", None)
    result.setdefault("enrichments", [])
    result.setdefault("summary", "")

    return result


# ─── 写入 wiki ────────────────────────────────────────────────────────────────

def commit_ingest(
    result: dict,
    wiki_path_root: Path,
    conflict_resolutions: dict[str, str],   # {conflict_id: "new"|"old"}
    approved_suggestions: list[str],         # 用户勾选的字段名列表
    source_dest_dir: Optional[Path] = None, # 将上传文件复制到此目录
    uploaded_file_path: Optional[Path] = None,
) -> dict:
    """
    将 LLM 解析结果写入 wiki 文件系统。
    返回 {"written": [...], "appended": [...], "skipped_conflicts": [...]}
    """
    written: list[str] = []
    appended: list[str] = []

    # 1. 应用冲突裁决（修改 frontmatter）
    for conflict in result.get("conflicts", []):
        cid = conflict.get("id", "")
        choice = conflict_resolutions.get(cid)
        field = conflict.get("field", "")
        if choice == "old" and field:
            # 保留旧值：从 frontmatter 中移除新值（LLM 可能已写入）
            result["frontmatter"].pop(field, None)
        # choice == "new" 或未裁决：保留 LLM 生成的值

    # 2. 写入主 wiki 文件
    target = wiki_path_root / result["wiki_path"]
    target.parent.mkdir(parents=True, exist_ok=True)

    # 构建文件内容：frontmatter + body
    post = _fm.Post(result["body"], **result["frontmatter"])
    file_content = _fm.dumps(post)

    # 如果是更新且有 history_block，追加在文件末尾
    if result.get("history_block"):
        file_content = file_content.rstrip() + "\n\n" + result["history_block"].strip() + "\n"

    target.write_text(file_content, encoding="utf-8")
    written.append(str(target.relative_to(wiki_path_root)))

    # 3. 关联实体追加
    for ru in result.get("related_updates", []):
        ru_path = wiki_path_root / ru["wiki_path"]
        if ru_path.exists():
            existing = ru_path.read_text(encoding="utf-8")
            if ru["append_text"] not in existing:
                ru_path.write_text(existing.rstrip() + "\n" + ru["append_text"] + "\n", encoding="utf-8")
                appended.append(str(ru_path.relative_to(wiki_path_root)))

    # 4. Schema 建议写入（只写用户勾选的）
    if approved_suggestions:
        suggestions_file = wiki_path_root / "schema-suggestions.md"
        approved_items = [
            s for s in result.get("schema_suggestions", [])
            if s.get("field") in approved_suggestions
        ]
        if approved_items:
            _append_schema_suggestions(suggestions_file, approved_items, result["wiki_path"])

    # 5. 未裁决的 Type B 冲突写入 pending-clarifications.md
    unresolved = [
        c for c in result.get("conflicts", [])
        if c.get("type") == "B" and c.get("id") not in conflict_resolutions
    ]
    if unresolved:
        _append_pending_clarifications(
            wiki_path_root / "pending-clarifications.md",
            unresolved,
            result["wiki_path"],
        )

    # 6. 写入 log.md
    log_file = wiki_path_root / "log.md"
    log_entry = result.get(
        "log_entry",
        f"[{date.today().isoformat()}] INGEST | {result.get('summary', result['wiki_path'])}",
    )
    _append_log(log_file, log_entry)

    # 7. 复制上传文件到 sources/proposals/（可选）
    if source_dest_dir and uploaded_file_path and uploaded_file_path.exists():
        source_dest_dir.mkdir(parents=True, exist_ok=True)
        dest = source_dest_dir / uploaded_file_path.name
        if not dest.exists():
            import shutil
            shutil.copy2(str(uploaded_file_path), str(dest))

    return {"written": written, "appended": appended}


# ─── 内部辅助函数 ─────────────────────────────────────────────────────────────

def _append_schema_suggestions(path: Path, items: list[dict], source_path: str) -> None:
    header = "| 建议字段名 | 示例值 | 来源 | 可信度 | 建议加入 Schema? |\n|---|---|---|---|---|\n"
    rows = "\n".join(
        f"| {i['field']} | {i.get('example_value','')} | {source_path} | {i.get('confidence','')} | ☐ |"
        for i in items
    )
    block = f"\n{rows}\n"

    if path.exists():
        content = path.read_text(encoding="utf-8")
        if "## 待确认字段建议" in content:
            # 找到表格末尾插入
            path.write_text(content.rstrip() + block, encoding="utf-8")
        else:
            path.write_text(content + f"\n## 待确认字段建议\n\n{header}" + block, encoding="utf-8")
    else:
        path.write_text(f"## 待确认字段建议\n\n{header}{block}", encoding="utf-8")


def _append_pending_clarifications(path: Path, conflicts: list[dict], source_path: str) -> None:
    lines = []
    for c in conflicts:
        lines.append(f"\n### [CONFLICT-WEB-{c.get('id','')}] {c.get('field','未知字段')}")
        lines.append(f"- 来源：{source_path}")
        lines.append(f"- 现有值：{c.get('existing_value', 'N/A')}")
        lines.append(f"- 新值：{c.get('new_value', 'N/A')}")
        lines.append(f"- 建议：{c.get('recommendation', '')}")
        lines.append("- 裁决：[ ] 以现有值为准  [ ] 以新值为准  [ ] 两者均保留\n")

    block = "\n".join(lines)
    if path.exists():
        path.write_text(path.read_text(encoding="utf-8").rstrip() + "\n" + block, encoding="utf-8")
    else:
        path.write_text(f"## 待澄清项\n{block}", encoding="utf-8")


def _append_log(path: Path, entry: str) -> None:
    if path.exists():
        path.write_text(path.read_text(encoding="utf-8").rstrip() + f"\n{entry}\n", encoding="utf-8")
    else:
        path.write_text(f"# Wiki 操作日志\n\n{entry}\n", encoding="utf-8")


def _append_section(file_path: Path, section_tag: str, content: str) -> None:
    """在 wiki 页面的指定 [section_tag] 分区末尾追加内容，分区不存在则在文件末尾创建。"""
    text = file_path.read_text(encoding="utf-8")
    stripped_content = content.strip()

    # 幂等检查：内容已存在则跳过
    if stripped_content and stripped_content[:80] in text:
        return

    pattern = re.compile(
        rf"(## \[{re.escape(section_tag)}\][^\n]*\n)(.*?)(?=\n## \[|\Z)",
        re.DOTALL,
    )
    m = pattern.search(text)
    if m:
        insert_pos = m.end(2)
        text = text[:insert_pos].rstrip() + "\n" + stripped_content + "\n" + text[insert_pos:]
    else:
        text = text.rstrip() + f"\n\n## [{section_tag}]\n{stripped_content}\n"

    file_path.write_text(text, encoding="utf-8")


# ─── Experience Ingest 专用函数 ───────────────────────────────────────────────

def load_experience_config(schema_dir: Path) -> dict:
    """
    从 domain-config.xlsx 读取 Experience Types 和 Experience Section Types Sheet。
    返回结构化 dict 供 build_experience_prompt 使用。
    """
    xlsx_path = schema_dir / "domain-config.xlsx"
    result: dict = {"experience_types": [], "section_types": []}
    if not xlsx_path.exists():
        return result

    try:
        import openpyxl
    except ImportError:
        return result

    wb = openpyxl.load_workbook(str(xlsx_path), read_only=True, data_only=True)

    for sheet in wb.worksheets:
        rows = [r for r in sheet.iter_rows(values_only=True) if any(c is not None for c in r)]
        if not rows:
            continue
        headers = [str(c).strip() if c else "" for c in rows[0]]

        if sheet.title == "Experience Types":
            for row in rows[1:]:
                record = {headers[i]: (str(row[i]).strip() if row[i] is not None else "") for i in range(len(headers))}
                result["experience_types"].append(record)

        elif sheet.title == "Experience Section Types":
            for row in rows[1:]:
                record = {headers[i]: (str(row[i]).strip() if row[i] is not None else "") for i in range(len(headers))}
                result["section_types"].append(record)

    wb.close()
    return result


def scan_experience_seq(wiki_path: Path, subtype_abbr: str) -> int:
    """
    扫描 wiki/experiences/{subtype}/ 目录，返回当前最大序号。
    若目录为空或不存在，返回 0（commit 时加 1 得到 01）。
    """
    exp_dir = wiki_path / "experiences" / subtype_abbr
    if not exp_dir.exists():
        return 0
    max_seq = 0
    for f in exp_dir.glob("*.md"):
        # 文件名格式：exp-{year}-{subtype}-{seq:02d}.md
        parts = f.stem.split("-")
        if len(parts) >= 4:
            try:
                seq = int(parts[-1])
                max_seq = max(max_seq, seq)
            except ValueError:
                pass
    return max_seq


_EXPERIENCE_SYSTEM_PROMPT = textwrap.dedent("""
你是企业知识库（LLM Wiki Kit）的经验摄入引擎。
你的任务是将运营活动报告（接待报告、月报、经验沉淀等）转化为 experience 实体，
并识别需要反哺到关联 wiki 实体的增益内容。

## 输出要求（严格遵守）
- 只输出一个合法的 JSON 对象，不添加任何说明文字
- JSON 必须放在 ```json ... ``` 代码块中
- 字符串中的换行用 \\n 表示

## experience 实体规范
- type 固定为 "experience"
- subtype 从 domain-config Experience Types 表中选择最匹配的 subtype_abbr
- event_date 取报告中活动发生的日期（非文件日期）
- entity_id 占位符填 "exp-{year}-{subtype}-TBD"，commit 时系统会替换为正确序号
- wiki_path 格式：experiences/{subtype}/exp-{year}-{subtype}-TBD.md
- dimensions 按 domain-config Experience Types 表的 coreDimensions 字段逐项提取
- subject 字段：匹配报告中提及的现有 wiki 实体路径（仅使用现有实体列表中存在的）

## 增益规则
- qa_append：从 [qa-record] 分区提取真实问答，target 为关联 modules/ 实体
  每条格式：Q: 问题\\nA: 回答（来源：{entity_id} / {event_date}）
- narration_append：从 [narration-record] 分区提取实际话术，target 为关联 modules/ 实体
  每条格式：> 话术内容（角色：讲解员，来源：{entity_id} / {event_date}）
- preference_append：从客户关注点提取偏好，target 为关联 clients/ 实体
  格式：- {event_date} {偏好描述}（来源：{entity_id}）
- scenarios_suggest：从 [visitor-insights] 发现新适用场景，target 为关联 modules/ 实体
  仅建议，不直接修改 scenarios 字段

## 自发现字段
对 Schema 未定义但有高价值的字段，放入 schema_suggestions 数组。
""").strip()


def build_experience_prompt(
    filename: str,
    extracted_text: str,
    domain_config: str,
    existing_entities: list[dict],
    experience_config: dict,
    today: str,
) -> tuple[str, str]:
    """返回 (system_prompt, user_prompt)，专为 experience 报告构建。"""

    if existing_entities:
        entity_lines = ["| wiki_path | title | type | status |", "|---|---|---|---|"]
        for e in existing_entities[:80]:
            entity_lines.append(
                f"| {e['wiki_path']} | {e['title']} | {e['type']} | {e['status']} |"
            )
        entities_section = "\n".join(entity_lines)
    else:
        entities_section = "（暂无现有实体）"

    # 格式化 experience 配置
    exp_types_lines = []
    for et in experience_config.get("experience_types", []):
        exp_types_lines.append(
            f"- subtype_abbr={et.get('subtype_abbr','')} | dimensions={et.get('coreDimensions','')} | sections={et.get('applicableSections','')}"
        )
    exp_types_section = "\n".join(exp_types_lines) if exp_types_lines else "（未配置）"

    section_types_lines = []
    for st in experience_config.get("section_types", []):
        section_types_lines.append(
            f"- [{st.get('sectionTag','')}] enrichTarget={st.get('enrichTarget','')} | appliesTo={st.get('applicableSubtype','')}"
        )
    section_types_section = "\n".join(section_types_lines) if section_types_lines else "（未配置）"

    user_prompt = textwrap.dedent(f"""
## 领域配置（domain-config.xlsx 完整内容）

{domain_config}

---

## Experience Types 配置

{exp_types_section}

## Experience Section Types 配置

{section_types_section}

---

## 现有知识库实体索引（用于 subject 字段匹配）

{entities_section}

---

## 待摄入文件

文件名：{filename}
摄入日期：{today}

内容：
{extracted_text}

---

## 输出 JSON 格式（严格按此结构）

```json
{{
  "action": "create",
  "wiki_path": "experiences/session/exp-{today[:4]}-session-TBD.md",
  "entity_id": "exp-{today[:4]}-session-TBD",
  "title": "YYYY-MM-DD 某客户 接待活动报告",
  "frontmatter": {{
    "title": "...",
    "type": "experience",
    "entity_id": "exp-{today[:4]}-session-TBD",
    "subtype": "session",
    "event_date": "YYYY-MM-DD",
    "status": "active",
    "created": "{today}",
    "updated": "{today}",
    "dimensions": {{}},
    "subject": [],
    "exhibited_assets": [],
    "sources": ["{filename}"]
  }},
  "body": "## [visitor-insights]\\n内容…\\n\\n## [qa-record]\\nQ: …\\nA: …\\n",
  "enrichments": [
    {{
      "enrich_type": "qa_append",
      "target": "modules/某模块.md",
      "content": "Q: 问题\\nA: 回答（来源：exp-id / date）",
      "source_label": "exp-id / date"
    }},
    {{
      "enrich_type": "narration_append",
      "target": "modules/某模块.md",
      "content": "> 话术内容（角色：讲解员，来源：exp-id / date）",
      "source_label": "exp-id / date"
    }},
    {{
      "enrich_type": "preference_append",
      "target": "clients/某客户.md",
      "content": "- YYYY-MM-DD 偏好描述（来源：exp-id）"
    }},
    {{
      "enrich_type": "scenarios_suggest",
      "target": "modules/某模块.md",
      "suggestion": "建议的 scenario 标签值",
      "reason": "来自 visitor-insights 的观察依据"
    }}
  ],
  "schema_suggestions": [],
  "summary": "一句话摘要"
}}
```

enrichments 数组只填实际发现的增益项，不填无内容的类型。
请严格按此格式输出，不添加其他文字。
""").strip()

    return _EXPERIENCE_SYSTEM_PROMPT, user_prompt


def commit_experience_ingest(
    result: dict,
    wiki_path_root: Path,
    approved_suggestions: list[str],
    source_dest_dir: Optional[Path] = None,
    uploaded_file_path: Optional[Path] = None,
) -> dict:
    """
    将 experience LLM 解析结果写入 wiki 文件系统。
    1. 写主 experience 实体文件（entity_id 序号由扫描确定）
    2. 执行跨实体增益（enrichments）
    3. 写 log.md（EXPERIENCE-INGEST + EXPERIENCE-ENRICH）
    """
    written: list[str] = []
    enriched: list[str] = []

    # 1. 确定正确的 entity_id 和 wiki_path（LLM 输出 TBD 占位，commit 时替换）
    fm = result["frontmatter"]
    subtype = fm.get("subtype", "session")
    event_year = (fm.get("event_date") or date.today().isoformat())[:4]
    next_seq = scan_experience_seq(wiki_path_root, subtype) + 1
    real_entity_id = f"exp-{event_year}-{subtype}-{next_seq:02d}"

    # 替换 TBD 占位符
    real_wiki_path = f"experiences/{subtype}/{real_entity_id}.md"
    fm["entity_id"] = real_entity_id
    fm["title"] = result.get("title", fm.get("title", real_entity_id))
    result["body"] = result.get("body", "").replace("exp-TBD", real_entity_id)

    # 2. 写主 experience 实体文件
    target = wiki_path_root / real_wiki_path
    target.parent.mkdir(parents=True, exist_ok=True)

    import frontmatter as _fm_mod
    post = _fm_mod.Post(result["body"], **fm)
    target.write_text(_fm_mod.dumps(post), encoding="utf-8")
    written.append(real_wiki_path)

    # 3. 执行跨实体增益
    enrich_counts: dict[str, int] = {"qa_append": 0, "narration_append": 0, "preference_append": 0, "scenarios_suggest": 0}
    for enrichment in result.get("enrichments", []):
        etype = enrichment.get("enrich_type", "")
        target_rel = enrichment.get("target", "")
        if not target_rel:
            continue
        target_path = wiki_path_root / target_rel

        if etype in ("qa_append", "narration_append"):
            if not target_path.exists():
                continue
            section = "qa" if etype == "qa_append" else "narration"
            _append_section(target_path, section, enrichment.get("content", ""))
            # version 递增
            _increment_version(target_path)
            enriched.append(target_rel)
            enrich_counts[etype] += 1

        elif etype == "preference_append":
            if not target_path.exists():
                continue
            content = enrichment.get("content", "").strip()
            if content:
                text = target_path.read_text(encoding="utf-8")
                if content not in text:
                    if "preference_notes" in text:
                        text = re.sub(
                            r"(preference_notes[^\n]*\n)",
                            lambda m: m.group(0) + content + "\n",
                            text, count=1,
                        )
                    else:
                        text = text.rstrip() + f"\n\n## preference_notes\n{content}\n"
                    target_path.write_text(text, encoding="utf-8")
                    _increment_version(target_path)
                    enriched.append(target_rel)
                    enrich_counts[etype] += 1

        elif etype == "scenarios_suggest":
            suggestion = enrichment.get("suggestion", "")
            reason = enrichment.get("reason", "")
            if suggestion:
                suggestions_file = wiki_path_root / "schema-suggestions.md"
                _append_schema_suggestions(
                    suggestions_file,
                    [{"field": f"scenarios:{suggestion}", "example_value": suggestion, "reason": reason, "confidence": "medium"}],
                    target_rel,
                )
                enrich_counts[etype] += 1

    # 4. Schema 建议（approved_suggestions 来自用户勾选）
    if approved_suggestions and result.get("schema_suggestions"):
        suggestions_file = wiki_path_root / "schema-suggestions.md"
        approved_items = [s for s in result["schema_suggestions"] if s.get("field") in approved_suggestions]
        if approved_items:
            _append_schema_suggestions(suggestions_file, approved_items, real_wiki_path)

    # 5. 复制上传文件到 sources/reports/
    if source_dest_dir and uploaded_file_path and uploaded_file_path.exists():
        source_dest_dir.mkdir(parents=True, exist_ok=True)
        dest = source_dest_dir / uploaded_file_path.name
        if not dest.exists():
            import shutil
            shutil.copy2(str(uploaded_file_path), str(dest))

    # 6. 写 log.md
    today_str = date.today().isoformat()
    log_file = wiki_path_root / "log.md"
    subject_list = ", ".join(fm.get("subject", []))
    _append_log(log_file, f"[{today_str}] EXPERIENCE-INGEST | {result.get('summary', real_entity_id)} | 新建: {real_wiki_path} | subject: {subject_list}")
    if any(v > 0 for v in enrich_counts.values()):
        _append_log(
            log_file,
            f"[{today_str}] EXPERIENCE-ENRICH | {real_entity_id} | "
            f"qa:{enrich_counts['qa_append']} narration:{enrich_counts['narration_append']} "
            f"preference:{enrich_counts['preference_append']} scenarios_suggest:{enrich_counts['scenarios_suggest']}",
        )

    return {"written": written, "enriched": enriched, "entity_id": real_entity_id}


def _increment_version(file_path: Path) -> None:
    """将 wiki 页面 frontmatter 中的 version 字段递增（v1→v2→v3）。"""
    try:
        import frontmatter as _fm_mod
        post = _fm_mod.loads(file_path.read_text(encoding="utf-8"))
        ver = str(post.metadata.get("version", "v1"))
        m = re.match(r"v(\d+)", ver)
        new_ver = f"v{int(m.group(1)) + 1}" if m else "v2"
        post.metadata["version"] = new_ver
        post.metadata["updated"] = date.today().isoformat()
        file_path.write_text(_fm_mod.dumps(post), encoding="utf-8")
    except Exception:
        pass
